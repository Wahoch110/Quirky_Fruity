"""Generate Week 5 Version Control slides (PPTX) with styling."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

# Palette: slate + sky (DevOps / pipeline feel)
C_BG_DARK = RGBColor(15, 23, 42)  # slate-900
C_BG_LIGHT = RGBColor(248, 250, 252)  # slate-50
C_RIBBON = RGBColor(30, 58, 95)  # deep blue
C_ACCENT = RGBColor(14, 165, 233)  # sky-500
C_ACCENT_SOFT = RGBColor(125, 211, 252)  # sky-300
C_TEXT = RGBColor(30, 41, 59)  # slate-800
C_TEXT_MUTED = RGBColor(100, 116, 139)  # slate-500
C_WHITE = RGBColor(255, 255, 255)


def set_notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def apply_title_slide_text(slide) -> None:
    title = slide.shapes.title
    tf = title.text_frame
    tf.word_wrap = True
    for para in tf.paragraphs:
        para.font.name = "Calibri"
        para.font.size = Pt(40)
        para.font.bold = True
        para.font.color.rgb = C_WHITE
        para.space_after = Pt(12)

    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            stf = shape.text_frame
            stf.word_wrap = True
            for para in stf.paragraphs:
                para.font.name = "Calibri"
                para.font.size = Pt(22)
                para.font.color.rgb = C_ACCENT_SOFT
                para.font.italic = True
            break


def add_title_slide(prs: Presentation, title: str, subtitle: str, notes: str) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            shape.text = subtitle
            break

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_BG_DARK

    strip_h = Inches(0.38)
    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        prs.slide_height - strip_h,
        prs.slide_width,
        strip_h,
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = C_ACCENT
    strip.line.fill.background()

    apply_title_slide_text(slide)

    # Decorative: small pipeline / node circles (left side accent)
    dot_y = Inches(2.2)
    for i, x in enumerate([Inches(0.85), Inches(1.15), Inches(1.45)]):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, dot_y + Inches(i * 0.08), Inches(0.14), Inches(0.14))
        dot.fill.solid()
        dot.fill.fore_color.rgb = C_ACCENT if i < 2 else C_ACCENT_SOFT
        dot.line.fill.background()

    set_notes(slide, notes)


def style_content_slide(slide, prs: Presentation) -> None:
    """Light background, full-width ribbon title, readable body."""
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_BG_LIGHT

    title_shape = slide.shapes.title
    title_shape.left = 0
    title_shape.top = 0
    title_shape.width = prs.slide_width
    title_shape.height = Inches(1.2)

    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = C_RIBBON
    title_shape.line.fill.background()

    ttf = title_shape.text_frame
    ttf.margin_left = Inches(0.5)
    ttf.margin_right = Inches(0.5)
    ttf.margin_top = Inches(0.35)
    ttf.word_wrap = True
    for para in ttf.paragraphs:
        para.font.name = "Calibri"
        para.font.size = Pt(28)
        para.font.bold = True
        para.font.color.rgb = C_WHITE

    # Left accent bar (vertical)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        Inches(1.2),
        Inches(0.12),
        prs.slide_height - Inches(1.2),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()

    # Content area placeholder
    ph = slide.placeholders[1]
    ph.left = Inches(0.55)
    ph.top = Inches(1.38)
    ph.width = prs.slide_width - Inches(1.1)
    ph.height = prs.slide_height - Inches(1.38) - Inches(0.45)


def style_body_paragraphs(body, bullet_size: int = 21) -> None:
    for para in body.paragraphs:
        para.font.name = "Calibri"
        para.font.size = Pt(bullet_size)
        para.font.color.rgb = C_TEXT
        para.space_after = Pt(10)
        para.line_spacing = 1.15
        # Bullet color via first run if needed — set on paragraph level
        if para.level == 0:
            para.font.color.rgb = C_TEXT
        else:
            para.font.color.rgb = C_TEXT_MUTED


def add_title_content_slide(prs: Presentation, title: str, bullets: list[str], notes: str) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    style_content_slide(slide, prs)

    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.level = 0
    style_body_paragraphs(body)

    set_notes(slide, notes)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Week 5: Version Control Fundamentals",
        'Building the "Source" of our Pipeline (Git & GitHub)',
        '"Welcome to the practical phase of the course! For the last four weeks, we talked about automated pipelines and cloud servers. But before we can automate our code, we need a secure, reliable place to store it. Today, we build the very first stage of the CI/CD pipeline: Source Control."',
    )

    add_title_content_slide(
        prs,
        'The "Final_v2_USE_THIS" Problem',
        [
            "How do you currently share group project code? USB drives? Emailing ZIP files?",
            "What happens when two people edit the exact same file at the same time?",
            "Result: Overwritten work, lost code, and total chaos.",
        ],
        '"We\'ve all been there. You name a file Project_Final, then Project_Final_v2, then Project_Real_Final_DO_NOT_TOUCH. This is a nightmare for professional software teams. We need a system that tracks exactly who changed what, and prevents team members from accidentally deleting each other\'s work."',
    )

    add_title_content_slide(
        prs,
        "What is Version Control? (The Time Machine)",
        [
            "A system that records changes to a file (or set of files) over time.",
            "The Time Machine: You can instantly recall specific, older versions of your project.",
            "The Sandbox: You can safely experiment without breaking the main application.",
        ],
        '"Version control is like a time machine for your code. If you make a massive mistake and the app crashes, you don\'t have to spend five hours hitting \'undo\' or trying to find the bug. You simply tell the system to roll back to the exact snapshot from yesterday when everything worked perfectly."',
    )

    add_title_content_slide(
        prs,
        "The Crucial Distinction: Git vs. GitHub",
        [
            "Git: The Engine. A command-line tool installed locally on your laptop that tracks file changes.",
            "GitHub: The Cloud. A website where you host your Git repositories so the whole world (or your team) can see them.",
        ],
        '"This is a guaranteed exam question and interview question. Do not mix these up. Git is the camera taking pictures of your code. GitHub is Instagram where you post the pictures for others to see. You can use Git without GitHub, but you cannot use GitHub without Git."',
    )

    add_title_content_slide(
        prs,
        "Core Vocabulary (Local)",
        [
            "Repository (Repo): The project folder that Git is actively tracking.",
            'Commit: A saved "snapshot" of your code at a specific point in time, always accompanied by a message explaining why you made the change.',
        ],
        '"When you type git commit, you are permanently saving the state of your files at that exact second. A project is just a long history of hundreds of commits. Your commit message should always be clear, like \'Added login button\', not \'Fixed stuff\'."',
    )

    add_title_content_slide(
        prs,
        "Core Vocabulary (Cloud)",
        [
            "Push: Uploading your local commits from your laptop up to the remote GitHub server.",
            "Pull: Downloading the latest commits from the remote GitHub server down to your laptop.",
        ],
        '"You will use these two commands every single day of your career. When you finish your work for the day, you push it to the cloud. When you start your work the next morning, you pull down whatever updates your teammates pushed while you were asleep."',
    )

    add_title_content_slide(
        prs,
        "The Standard Developer Workflow",
        [
            "Write Code → 2. Save File → 3. git commit (Snapshot) → 4. git push (Upload).",
        ],
        '"This is the heartbeat of DevOps. Once we learn this rhythm today, we will eventually connect step 4 to our CI/CD pipeline. Every time you \'push\', the automated factory will wake up and start testing your code."',
    )

    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Lab Time!"
    style_content_slide(slide, prs)

    body = slide.placeholders[1].text_frame
    body.clear()
    lines = [
        ("Today's Goal:", 0),
        ("Configure Git on your machine.", 1),
        ("Create a local repository.", 1),
        ("Make your first commit.", 1),
        ("Push your code to the internet.", 1),
    ]
    for i, (text, level) in enumerate(lines):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = text
        p.level = level
    style_body_paragraphs(body)

    set_notes(
        slide,
        '"Enough theory. Open up your terminals, log into GitHub, and let\'s get our hands dirty."',
    )

    out = Path(__file__).resolve().parent.parent / "Week5_Version_Control_Fundamentals.pptx"
    prs.save(str(out))
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
